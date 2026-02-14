#!/usr/bin/env python3
"""Generate presentation.pptx from the HTML presentation content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Constants ───────────────────────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(BASE, "slide-images")
LOGO = os.path.join(BASE, "MarjUnterbergNursHealthStud_Logo_Vert_White.png")

NAVY    = RGBColor(0x0A, 0x16, 0x28)
BLUE    = RGBColor(0x1E, 0x3A, 0x5F)
TEAL    = RGBColor(0x2A, 0x9D, 0x8F)
GOLD    = RGBColor(0xE9, 0xC4, 0x6A)
CORAL   = RGBColor(0xE7, 0x6F, 0x51)
WHITE   = RGBColor(0xF8, 0xF9, 0xFA)
LIGHT   = RGBColor(0xE8, 0xEC, 0xF1)
DIM     = RGBColor(0x99, 0x99, 0x99)
BLUE_ACC = RGBColor(0x4A, 0x90, 0xD9)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_TITLE = "Georgia"       # fallback for Playfair Display
FONT_BODY  = "Calibri"       # fallback for Inter


# ─── Helpers ─────────────────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_gradient_bg(slide, c1, c2):
    """Two-stop gradient background."""
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_stops[1].position = 1.0


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, font_name=FONT_BODY, size=Pt(18), color=WHITE, bold=False, italic=False, alignment=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return p


def add_run(paragraph, text, font_name=FONT_BODY, size=Pt(18), color=WHITE, bold=False, italic=False):
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return run


def add_paragraph(tf, text="", font_name=FONT_BODY, size=Pt(18), color=WHITE, bold=False, italic=False, alignment=PP_ALIGN.LEFT, space_before=Pt(0), space_after=Pt(0)):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    if text:
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
    return p


def add_accent_line(slide, left, top, width=Inches(1.2), height=Pt(4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape


def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        kwargs = {"left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        return slide.shapes.add_picture(path, **kwargs)
    return None


def add_card_bg(slide, left, top, width, height, border_color=None):
    """Add a rounded rectangle card background."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x14, 0x20, 0x35)
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.color.rgb = RGBColor(0x2A, 0x35, 0x50)
        shape.line.width = Pt(1)
    return shape


def add_slide_number(slide, num, total=17):
    tb = add_textbox(slide, Inches(12.3), Inches(7.0), Inches(1.0), Inches(0.4))
    set_text(tb.text_frame, f"{num} / {total}", size=Pt(11), color=DIM, alignment=PP_ALIGN.RIGHT)


# ─── Slide Builders ──────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Title slide: Building Leaders from the Ground Up"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_gradient_bg(slide, NAVY, TEAL)

    # Badge
    tb = add_textbox(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.5))
    set_text(tb.text_frame, "● Monmouth University Scholarship Week 2026", size=Pt(16), color=GOLD, bold=True)

    # Title
    tb = add_textbox(slide, Inches(0.8), Inches(1.2), Inches(6.5), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_run(p, "Building\n", FONT_TITLE, Pt(52), WHITE, bold=True)
    add_run(p, "Leaders\n", FONT_TITLE, Pt(52), GOLD, bold=True)
    add_run(p, "from the ", FONT_TITLE, Pt(36), DIM, bold=True)
    add_run(p, "Ground Up", FONT_TITLE, Pt(44), TEAL, bold=True)

    # Accent line
    add_accent_line(slide, Inches(0.8), Inches(4.4), Inches(2))

    # Subtitle
    tb = add_textbox(slide, Inches(0.8), Inches(4.7), Inches(6), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_run(p, "How a Graduate Assistant Designed, Launched,\nand Grew the ", FONT_BODY, Pt(18), LIGHT)
    add_run(p, "SNHS Student\nAmbassador Program", FONT_BODY, Pt(18), GOLD, bold=True)

    # Presenter card
    card = add_card_bg(slide, Inches(7.8), Inches(0.8), Inches(3.3), Inches(4.5))

    tb = add_textbox(slide, Inches(7.9), Inches(1.0), Inches(3.1), Inches(0.3))
    set_text(tb.text_frame, "PRESENTED BY", size=Pt(10), color=DIM, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(7.9), Inches(1.4), Inches(3.1), Inches(0.6))
    set_text(tb.text_frame, "Bingjun Li", FONT_TITLE, Pt(32), WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(7.9), Inches(2.0), Inches(3.1), Inches(0.4))
    set_text(tb.text_frame, "M.S.Ed.", size=Pt(16), color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(7.9), Inches(2.5), Inches(3.1), Inches(0.4))
    set_text(tb.text_frame, "Graduate Assistant", size=Pt(16), color=LIGHT, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(7.9), Inches(3.2), Inches(3.1), Inches(0.3))
    set_text(tb.text_frame, "— MENTORED BY —", size=Pt(9), color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(7.9), Inches(3.5), Inches(3.1), Inches(0.5))
    set_text(tb.text_frame, "Dr. Clifford", FONT_TITLE, Pt(24), GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(7.9), Inches(4.0), Inches(3.1), Inches(0.4))
    set_text(tb.text_frame, "Acting Dean, SNHS", size=Pt(13), color=DIM, alignment=PP_ALIGN.CENTER)

    # Logo
    add_image_safe(slide, LOGO, Inches(11.4), Inches(1.2), height=Inches(3.8))

    add_slide_number(slide, 1)


def slide_02_hook(prs):
    """The Hook: shared vision"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # Left text
    tb = add_textbox(slide, Inches(0.8), Inches(1.0), Inches(5.5), Inches(0.6))
    set_text(tb.text_frame, "It started with a conversation", FONT_TITLE, Pt(24), DIM)

    tb = add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_run(p, "A ", FONT_TITLE, Pt(36), WHITE, bold=True)
    add_run(p, "shared vision", FONT_TITLE, Pt(36), GOLD, bold=True)
    add_run(p, ",\na mentor, and a plan", FONT_TITLE, Pt(36), WHITE, bold=True)

    # Quote block
    quote_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.3), Inches(5.5), Inches(1.5))
    quote_bg.fill.solid()
    quote_bg.fill.fore_color.rgb = RGBColor(0x1A, 0x25, 0x38)
    quote_bg.line.fill.background()
    # Gold left border
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.3), Pt(4), Inches(1.5))
    border.fill.solid()
    border.fill.fore_color.rgb = GOLD
    border.line.fill.background()

    tb = add_textbox(slide, Inches(1.0), Inches(3.4), Inches(5.2), Inches(1.3))
    set_text(tb.text_frame, "Dr. Clifford shared her ideas for a student leadership program. We brainstormed events, surveyed students, and shaped the workshops together based on real feedback.", size=Pt(17), color=WHITE, italic=True)

    tb = add_textbox(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(0.5))
    set_text(tb.text_frame, "With her guidance, I built it from the ground up.", size=Pt(17), color=LIGHT)

    # Right image
    add_image_safe(slide, os.path.join(IMG, "s2-hook.png"), Inches(7.0), Inches(0.5), height=Inches(6.3))

    add_slide_number(slide, 2)


def slide_03_who_am_i(prs):
    """Who Am I: From Classroom to Program Builder"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BLUE)

    tb = add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5.5), Inches(0.4))
    set_text(tb.text_frame, "Who Am I", size=Pt(20), color=GOLD, bold=True)

    tb = add_textbox(slide, Inches(0.8), Inches(1.0), Inches(5.5), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_run(p, "From Classroom\nto ", FONT_TITLE, Pt(36), WHITE, bold=True)
    add_run(p, "Program Builder", FONT_TITLE, Pt(36), GOLD, bold=True)

    add_accent_line(slide, Inches(0.8), Inches(2.4))

    bullets = [
        ("Bingjun Li", " — M.S.Ed. Graduate Assistant"),
        ("", "Assigned to SNHS under Dr. Clifford"),
        ("", "Tasked with building the Student Ambassador Program from concept to reality"),
        ("", "Education background meeting healthcare leadership — a unique intersection"),
    ]
    tb = add_textbox(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (bold_part, rest) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        add_run(p, "● ", FONT_BODY, Pt(17), GOLD)
        if bold_part:
            add_run(p, bold_part, FONT_BODY, Pt(17), WHITE, bold=True)
        add_run(p, rest, FONT_BODY, Pt(17), LIGHT)

    # Right side: intersection visual (static)
    # Education circle
    edu = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(1.5), Inches(2.5), Inches(2.5))
    edu.fill.solid()
    edu.fill.fore_color.rgb = RGBColor(0x1A, 0x30, 0x50)
    edu.line.color.rgb = GOLD
    edu.line.width = Pt(2)
    tb = add_textbox(slide, Inches(7.7), Inches(2.2), Inches(2.1), Inches(1.2))
    tf = tb.text_frame
    set_text(tf, "🎓 Education", FONT_BODY, Pt(18), GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, "Curriculum · Pedagogy\nAssessment", size=Pt(12), color=LIGHT, alignment=PP_ALIGN.CENTER)

    # Healthcare circle
    hc = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.8), Inches(1.5), Inches(2.5), Inches(2.5))
    hc.fill.solid()
    hc.fill.fore_color.rgb = RGBColor(0x14, 0x30, 0x35)
    hc.line.color.rgb = TEAL
    hc.line.width = Pt(2)
    tb = add_textbox(slide, Inches(10.0), Inches(2.2), Inches(2.1), Inches(1.2))
    tf = tb.text_frame
    set_text(tf, "🩺 Healthcare", FONT_BODY, Pt(18), TEAL, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, "Leadership · Clinical\nCommunity", size=Pt(12), color=LIGHT, alignment=PP_ALIGN.CENTER)

    # Center label
    center_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.8), Inches(4.5), Inches(2.2), Inches(2.2))
    center_shape.fill.solid()
    center_shape.fill.fore_color.rgb = RGBColor(0x18, 0x2E, 0x42)
    center_shape.line.color.rgb = CORAL
    center_shape.line.width = Pt(2)
    tb = add_textbox(slide, Inches(8.9), Inches(4.9), Inches(2.0), Inches(1.4))
    tf = tb.text_frame
    set_text(tf, "Student\nAmbassador\nProgram", FONT_BODY, Pt(14), WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(7.5), Inches(6.9), Inches(4.8), Inches(0.4))
    set_text(tb.text_frame, "Where education meets healthcare leadership", size=Pt(13), color=DIM, italic=True, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 3)


def slide_04_mission(prs):
    """The Mission: What Is the Student Ambassador Program?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, NAVY, RGBColor(0x1A, 0x4A, 0x42))

    tb = add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.4))
    set_text(tb.text_frame, "The Mission", size=Pt(20), color=GOLD, bold=True)

    tb = add_textbox(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8))
    set_text(tb.text_frame, "What Is the Student Ambassador Program?", FONT_TITLE, Pt(36), WHITE, bold=True)

    add_accent_line(slide, Inches(0.8), Inches(2.0))

    tb = add_textbox(slide, Inches(0.8), Inches(2.3), Inches(10), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_run(p, "A ", FONT_BODY, Pt(18), LIGHT)
    add_run(p, "one-year leadership development program", FONT_BODY, Pt(18), GOLD, bold=True)
    add_run(p, " for 20-30 SNHS undergraduate students who represent the school, lead health promotion initiatives, and grow as future healthcare leaders.", FONT_BODY, Pt(18), LIGHT)

    # Tags
    tags = [("Transformational", GOLD), ("Democratic", TEAL), ("Adaptive", CORAL), ("Collaborative", GOLD)]
    left = Inches(0.8)
    for tag_text, tag_color in tags:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.5), Inches(2.2), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(tag_color[0] // 4, tag_color[1] // 4, tag_color[2] // 4)
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = tag_text
        run.font.size = Pt(15)
        run.font.color.rgb = tag_color
        run.font.bold = True
        left += Inches(2.4)

    # Four pillars
    pillars = [
        ("🎯", "Leadership\nTraining", "Workshops, styles,\ncase studies", GOLD),
        ("🤝", "Community\nOutreach", "Healthy Futures,\nlocal schools", TEAL),
        ("📈", "Professional\nDevelopment", "Public speaking,\nmentorship", CORAL),
        ("🩺", "Health Career\nExposure", "Simulations,\nfield trips", BLUE_ACC),
    ]
    for i, (icon, label, sub, color) in enumerate(pillars):
        x = Inches(0.8 + i * 3.1)
        # Bar accent
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(4.5), Inches(2.8), Pt(4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        # Card bg
        add_card_bg(slide, x, Inches(4.7), Inches(2.8), Inches(2.5), color)
        # Icon
        tb = add_textbox(slide, x, Inches(4.8), Inches(2.8), Inches(0.5))
        set_text(tb.text_frame, icon, size=Pt(28), alignment=PP_ALIGN.CENTER)
        # Label
        tb = add_textbox(slide, x, Inches(5.3), Inches(2.8), Inches(0.8))
        set_text(tb.text_frame, label, FONT_BODY, Pt(16), color, bold=True, alignment=PP_ALIGN.CENTER)
        # Sub
        tb = add_textbox(slide, x, Inches(6.1), Inches(2.8), Inches(0.8))
        set_text(tb.text_frame, sub, FONT_BODY, Pt(13), LIGHT, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 4)


def slide_05_what_i_built(prs):
    """What I Built: Designing Everything from Scratch"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    tb = add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.4))
    set_text(tb.text_frame, "What I Built", size=Pt(20), color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, "Designing ", FONT_TITLE, Pt(34), WHITE, bold=True)
    add_run(p, "Everything", FONT_TITLE, Pt(34), GOLD, bold=True)
    add_run(p, " from Scratch", FONT_TITLE, Pt(34), WHITE, bold=True)

    cards = [
        ("s5-curriculum.png", "Leadership Curriculum", "Transformational, democratic, adaptive & collaborative leadership workshops with case studies", GOLD),
        ("s5-assessment.png", "Assessment Tools", "Leadership questionnaire & activity interest survey to measure growth and guide programming", TEAL),
        ("s5-events.png", "Events & Workshops", "2 leadership workshops, 2 ice breaker events, plus full Spring 2026 programming calendar", CORAL),
        ("s5-outreach.png", "Outreach Materials", "Healthy Futures flyer, program poster, promotional materials for recruitment & community partners", BLUE_ACC),
    ]
    for i, (img_file, title, desc, color) in enumerate(cards):
        x = Inches(0.5 + i * 3.15)
        # Card background
        add_card_bg(slide, x, Inches(1.8), Inches(2.9), Inches(5.3), color)
        # Image
        add_image_safe(slide, os.path.join(IMG, img_file), x + Inches(0.1), Inches(1.9), width=Inches(2.7))
        # Title
        tb = add_textbox(slide, x + Inches(0.15), Inches(4.7), Inches(2.6), Inches(0.5))
        set_text(tb.text_frame, title, FONT_BODY, Pt(16), color, bold=True)
        # Description
        tb = add_textbox(slide, x + Inches(0.15), Inches(5.2), Inches(2.6), Inches(1.5))
        set_text(tb.text_frame, desc, size=Pt(13), color=LIGHT)

    add_slide_number(slide, 5)


def slide_06_timeline(prs):
    """Program Timeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BLUE)

    tb = add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.4))
    set_text(tb.text_frame, "The Journey", size=Pt(20), color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, "Program ", FONT_TITLE, Pt(34), WHITE, bold=True)
    add_run(p, "Timeline", FONT_TITLE, Pt(34), GOLD, bold=True)

    # Timeline track
    track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.5), Pt(6))
    track.fill.solid()
    track.fill.fore_color.rgb = TEAL
    track.line.fill.background()

    phases = [
        ("s6-foundation.png", "Foundation", "Fall 2025", "Grant Secured\nProgram Design\nRecruitment", GOLD),
        ("s6-launch.png", "Launch", "Nov 2025", "Ice Breaker Events\nTeam Building\n20-30 Ambassadors", TEAL),
        ("s6-workshops.png", "Workshops", "Nov-Dec 2025", "Leadership Training\nTransformational\nLeadership Focus", CORAL),
        ("s6-growth.png", "Growth", "Spring 2026", "Simulation Field Trip\nMentorship Training\nInterest Surveys", BLUE_ACC),
        ("s6-future.png", "Future", "2026+", "Community Outreach\nLeaders Symposium\nHealthy Futures", WHITE),
    ]
    for i, (img_file, phase, date, detail, color) in enumerate(phases):
        x = Inches(0.5 + i * 2.5)
        # Dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), Inches(1.65), Inches(0.25), Inches(0.25))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        # Image
        add_image_safe(slide, os.path.join(IMG, img_file), x + Inches(0.2), Inches(2.1), width=Inches(1.6))
        # Phase label
        tb = add_textbox(slide, x, Inches(4.2), Inches(2.2), Inches(0.4))
        set_text(tb.text_frame, phase, FONT_BODY, Pt(16), color, bold=True, alignment=PP_ALIGN.CENTER)
        # Date
        tb = add_textbox(slide, x, Inches(4.6), Inches(2.2), Inches(0.3))
        set_text(tb.text_frame, date, size=Pt(13), color=DIM, alignment=PP_ALIGN.CENTER)
        # Detail
        tb = add_textbox(slide, x, Inches(5.0), Inches(2.2), Inches(1.5))
        set_text(tb.text_frame, detail, size=Pt(13), color=LIGHT, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 6)


def slide_07_workshop1(prs):
    """Workshop 1: Leadership & Communication Skills"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    tb = add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4))
    set_text(tb.text_frame, "Workshop 1, November 18, 2025", size=Pt(18), color=GOLD, bold=True)

    tb = add_textbox(slide, Inches(0.8), Inches(1.0), Inches(5.5), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_run(p, "Leadership &\n", FONT_TITLE, Pt(34), WHITE, bold=True)
    add_run(p, "Communication Skills", FONT_TITLE, Pt(34), GOLD, bold=True)

    add_accent_line(slide, Inches(0.8), Inches(2.5))

    tb = add_textbox(slide, Inches(0.8), Inches(2.8), Inches(5), Inches(2.5))
    set_text(tb.text_frame,
        "My first workshop ever, designed from scratch. Introduced different leadership styles through real world case studies, group discussions, and interactive activities to build communication skills.",
        size=Pt(17), color=LIGHT)

    # Workshop images 2x2
    ws1_imgs = ["ws1-01.png", "ws1-03.png", "ws1-05.png", "ws1-09.png"]
    positions = [(Inches(7.0), Inches(0.5)), (Inches(10.0), Inches(0.5)),
                 (Inches(7.0), Inches(3.7)), (Inches(10.0), Inches(3.7))]
    for img_f, (x, y) in zip(ws1_imgs, positions):
        add_image_safe(slide, os.path.join(IMG, img_f), x, y, width=Inches(2.8))

    tb = add_textbox(slide, Inches(7.0), Inches(7.0), Inches(5.8), Inches(0.3))
    set_text(tb.text_frame, "Actual slides from Workshop 1", size=Pt(11), color=DIM, italic=True, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 7)


def slide_08_evolution(prs):
    """What I Learned & Changed"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, NAVY, RGBColor(0x1A, 0x4A, 0x42))

    tb = add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.4))
    set_text(tb.text_frame, "My Growth as a Facilitator", size=Pt(20), color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, "What I ", FONT_TITLE, Pt(34), WHITE, bold=True)
    add_run(p, "Learned", FONT_TITLE, Pt(34), GOLD, bold=True)
    add_run(p, " & Changed", FONT_TITLE, Pt(34), WHITE, bold=True)

    # Left column: What I noticed
    tb = add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4))
    set_text(tb.text_frame, "After Workshop 1, I noticed...", FONT_BODY, Pt(18), CORAL, bold=True)

    noticed = [
        "📝 Too much content delivery, not enough student interaction",
        "💬 Students wanted more discussion and real conversation",
        "⏱ Pacing felt rushed, needed to slow down and listen",
    ]
    for i, item in enumerate(noticed):
        tb = add_textbox(slide, Inches(0.8), Inches(2.3 + i * 0.9), Inches(5.5), Inches(0.8))
        set_text(tb.text_frame, item, size=Pt(16), color=LIGHT)

    # Right column: What I changed
    tb = add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.4))
    set_text(tb.text_frame, "So for Workshop 2, I...", FONT_BODY, Pt(18), TEAL, bold=True)

    changed = [
        "📊 Added a post-survey to capture feedback in real time",
        "🤝 Shifted to discussion-driven format, less lecturing, more facilitating",
        "🎓 Co-facilitated with Dr. Clifford, learned from watching her lead",
    ]
    for i, item in enumerate(changed):
        tb = add_textbox(slide, Inches(7.0), Inches(2.3 + i * 0.9), Inches(5.5), Inches(0.8))
        set_text(tb.text_frame, item, size=Pt(16), color=LIGHT)

    # Evolution arrow bar
    bar_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(5.5), Inches(10), Pt(6))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = TEAL
    bar_bg.line.fill.background()

    tb = add_textbox(slide, Inches(0.8), Inches(5.7), Inches(3), Inches(0.4))
    set_text(tb.text_frame, "Content Heavy", size=Pt(14), color=CORAL, bold=True)

    tb = add_textbox(slide, Inches(4.5), Inches(5.7), Inches(4), Inches(0.4))
    set_text(tb.text_frame, "Listen first, then redesign", size=Pt(14), color=DIM, italic=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(9.5), Inches(5.7), Inches(3), Inches(0.4))
    set_text(tb.text_frame, "Student Centered", size=Pt(14), color=TEAL, bold=True, alignment=PP_ALIGN.RIGHT)

    add_slide_number(slide, 8)


def slide_09_workshop2(prs):
    """Workshop 2: Transformational Leadership"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    tb = add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4))
    set_text(tb.text_frame, "Workshop 2, December 2, 2025", size=Pt(18), color=GOLD, bold=True)

    tb = add_textbox(slide, Inches(0.8), Inches(1.0), Inches(5.5), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_run(p, "Transformational\n", FONT_TITLE, Pt(34), WHITE, bold=True)
    add_run(p, "Leadership", FONT_TITLE, Pt(34), GOLD, bold=True)

    add_accent_line(slide, Inches(0.8), Inches(2.5))

    tb = add_textbox(slide, Inches(0.8), Inches(2.8), Inches(5), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_run(p, "A redesigned, discussion-driven workshop. Started with Simon Sinek's ", FONT_BODY, Pt(17), LIGHT)
    add_run(p, '"Start With Why"', FONT_BODY, Pt(17), GOLD, bold=True)
    add_run(p, ", asked students to reflect on their own why, and explored team dynamics through Patrick Lencioni's ", FONT_BODY, Pt(17), LIGHT)
    add_run(p, "Five Dysfunctions of a Team", FONT_BODY, Pt(17), WHITE, bold=True)
    add_run(p, ".", FONT_BODY, Pt(17), LIGHT)

    # Workshop images 2x2
    ws2_imgs = ["ws2-01.png", "ws2-03.png", "ws2-04.png", "ws2-11.png"]
    positions = [(Inches(7.0), Inches(0.5)), (Inches(10.0), Inches(0.5)),
                 (Inches(7.0), Inches(3.7)), (Inches(10.0), Inches(3.7))]
    for img_f, (x, y) in zip(ws2_imgs, positions):
        add_image_safe(slide, os.path.join(IMG, img_f), x, y, width=Inches(2.8))

    tb = add_textbox(slide, Inches(7.0), Inches(7.0), Inches(5.8), Inches(0.3))
    set_text(tb.text_frame, "Actual slides from Workshop 2", size=Pt(11), color=DIM, italic=True, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 9)


def slide_10_feedback(prs):
    """Student Feedback"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BLUE)

    tb = add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.4))
    set_text(tb.text_frame, "In Their Own Words", size=Pt(20), color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, "What the ", FONT_TITLE, Pt(34), WHITE, bold=True)
    add_run(p, "Ambassadors", FONT_TITLE, Pt(34), GOLD, bold=True)
    add_run(p, " Said", FONT_TITLE, Pt(34), WHITE, bold=True)

    # Quote cards
    quotes = [
        ('"B did great! It was an effective way of hearing how I could actually be a good leader. Good to see a little bit of Dr. Clifford\'s personality too."', "Post-workshop survey response", GOLD),
        ('"I really loved this activity that we did with the nursing ambassadors. I can\'t wait to see what future events has in store for us."', "Post-workshop survey response", TEAL),
        ('"I would like to get a better communication style when working with a team... gain the confidence in public speaking and talking with families, faculty, potential students."', "What students want from the program", CORAL),
    ]
    for i, (quote, attrib, color) in enumerate(quotes):
        y = Inches(1.6 + i * 1.8)
        # Card bg
        card = add_card_bg(slide, Inches(0.6), y, Inches(6.2), Inches(1.6), color)
        # Quote text
        tb = add_textbox(slide, Inches(0.8), y + Inches(0.1), Inches(5.8), Inches(1.1))
        set_text(tb.text_frame, quote, size=Pt(14), color=WHITE, italic=True)
        # Attribution
        tb = add_textbox(slide, Inches(0.8), y + Inches(1.2), Inches(5.8), Inches(0.3))
        set_text(tb.text_frame, f"— {attrib}", size=Pt(12), color=color)

    # Survey image
    add_image_safe(slide, os.path.join(IMG, "survey-3.png"), Inches(7.2), Inches(1.8), width=Inches(5.5))

    tb = add_textbox(slide, Inches(7.2), Inches(5.5), Inches(5.5), Inches(0.3))
    set_text(tb.text_frame, "Activity Interest Survey Results (3 responses)", size=Pt(11), color=DIM, italic=True, alignment=PP_ALIGN.CENTER)

    # Highlight box
    highlight_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(6.0), Inches(5.0), Inches(1.0))
    highlight_bg.fill.solid()
    highlight_bg.fill.fore_color.rgb = RGBColor(0x1E, 0x2A, 0x3F)
    highlight_bg.line.color.rgb = GOLD
    highlight_bg.line.width = Pt(1)

    tb = add_textbox(slide, Inches(7.6), Inches(6.1), Inches(4.8), Inches(0.4))
    set_text(tb.text_frame, "Leadership & Communication: 100% interest", size=Pt(15), color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(7.6), Inches(6.5), Inches(4.8), Inches(0.3))
    set_text(tb.text_frame, "Students want practical, interactive experiences", size=Pt(13), color=DIM, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 10)


def slide_11_challenges(prs):
    """Challenges"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, BLUE, RGBColor(0x4A, 0x2C, 0x2A))

    # Left side
    tb = add_textbox(slide, Inches(0.5), Inches(0.5), Inches(5.5), Inches(0.4))
    set_text(tb.text_frame, "Real Talk", size=Pt(20), color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(0.5), Inches(1.0), Inches(5.5), Inches(1.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, "The ", FONT_TITLE, Pt(34), WHITE, bold=True)
    add_run(p, "Challenges", FONT_TITLE, Pt(34), CORAL, bold=True)
    add_run(p, "\nNobody Warns You About", FONT_TITLE, Pt(34), WHITE, bold=True)

    add_image_safe(slide, os.path.join(IMG, "s11-challenges.png"), Inches(1.0), Inches(2.8), height=Inches(4.0))

    # Right side: challenge cards
    challenges = [
        ("Logistics & Coordination", "Scheduling 20-30 busy undergrads, booking rooms, managing budgets, handling permissions — the invisible work that makes everything else possible.", CORAL),
        ("Curriculum Design", "How do you teach leadership to nursing and health students in a way that's practical, not theoretical? Finding the right balance was an ongoing experiment.", GOLD),
        ("Topic Selection & Engagement", "Which leadership styles matter most? How do you keep participation high? Every choice had to be intentional and student-centered.", TEAL),
    ]
    for i, (title, desc, color) in enumerate(challenges):
        y = Inches(0.8 + i * 2.2)
        card = add_card_bg(slide, Inches(6.8), y, Inches(5.8), Inches(2.0), color)
        tb = add_textbox(slide, Inches(7.0), y + Inches(0.15), Inches(5.4), Inches(0.4))
        set_text(tb.text_frame, title, FONT_BODY, Pt(16), color, bold=True)
        tb = add_textbox(slide, Inches(7.0), y + Inches(0.6), Inches(5.4), Inches(1.2))
        set_text(tb.text_frame, desc, size=Pt(14), color=LIGHT)

    add_slide_number(slide, 11)


def slide_12_breakthrough(prs):
    """Breakthrough Moment"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # Left side
    tb = add_textbox(slide, Inches(0.5), Inches(0.5), Inches(5.5), Inches(0.4))
    set_text(tb.text_frame, "The Turning Point", size=Pt(20), color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(0.5), Inches(1.0), Inches(5.5), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, "When It All ", FONT_TITLE, Pt(34), WHITE, bold=True)
    add_run(p, "Clicked", FONT_TITLE, Pt(34), GOLD, bold=True)

    add_image_safe(slide, os.path.join(IMG, "s12-breakthrough.png"), Inches(1.0), Inches(2.2), height=Inches(4.5))

    # Right side
    # Quote
    quote_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(0.6), Inches(5.8), Inches(1.6))
    quote_bg.fill.solid()
    quote_bg.fill.fore_color.rgb = RGBColor(0x1A, 0x25, 0x38)
    quote_bg.line.fill.background()
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(0.6), Pt(4), Inches(1.6))
    border.fill.solid()
    border.fill.fore_color.rgb = GOLD
    border.line.fill.background()
    tb = add_textbox(slide, Inches(7.0), Inches(0.7), Inches(5.4), Inches(1.4))
    set_text(tb.text_frame, "The students were engaged from the start, in both workshops. What really changed was me. By the second workshop, I stopped anticipating their answers. I stopped planning how to respond before they finished speaking.", size=Pt(15), color=WHITE, italic=True)

    tb = add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.8), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_run(p, "I learned to ", FONT_BODY, Pt(16), LIGHT)
    add_run(p, "truly listen", FONT_BODY, Pt(16), GOLD, bold=True)
    add_run(p, ", to be present in the moment, and to let the conversation flow naturally instead of controlling it.", FONT_BODY, Pt(16), LIGHT)

    tb = add_textbox(slide, Inches(6.8), Inches(3.7), Inches(5.8), Inches(1.0))
    set_text(tb.text_frame, "That shift made all the difference. I became more engaged, more curious, and more connected to the students I was serving.", size=Pt(16), color=LIGHT)

    # Bottom images
    add_image_safe(slide, os.path.join(IMG, "ws2-09.png"), Inches(6.8), Inches(4.9), width=Inches(2.8))
    add_image_safe(slide, os.path.join(IMG, "ws2-10.png"), Inches(9.8), Inches(4.9), width=Inches(2.8))

    tb = add_textbox(slide, Inches(6.8), Inches(7.0), Inches(5.8), Inches(0.3))
    set_text(tb.text_frame, "Case study slides that sparked real debate", size=Pt(11), color=DIM, italic=True, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 12)


def slide_13_impact(prs):
    """Impact & Numbers"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, NAVY, RGBColor(0x1A, 0x4A, 0x42))

    tb = add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.4))
    set_text(tb.text_frame, "Impact So Far", size=Pt(20), color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, "By the ", FONT_TITLE, Pt(34), WHITE, bold=True)
    add_run(p, "Numbers", FONT_TITLE, Pt(34), GOLD, bold=True)

    stats = [
        ("20", "Undergraduate\nAmbassadors", GOLD),
        ("4", "Events\nCompleted", TEAL),
        ("4", "Leadership Styles\nTaught", CORAL),
        ("$10K", "Diversity Innovation\nGrant Secured", BLUE_ACC),
    ]
    for i, (number, label, color) in enumerate(stats):
        x = Inches(0.8 + i * 3.1)
        # Card bg
        card = add_card_bg(slide, x, Inches(1.8), Inches(2.8), Inches(3.5), color)
        # Number
        tb = add_textbox(slide, x, Inches(2.0), Inches(2.8), Inches(1.5))
        set_text(tb.text_frame, number, FONT_TITLE, Pt(64), color, bold=True, alignment=PP_ALIGN.CENTER)
        # Label
        tb = add_textbox(slide, x, Inches(3.5), Inches(2.8), Inches(1.0))
        set_text(tb.text_frame, label, FONT_BODY, Pt(16), LIGHT, bold=True, alignment=PP_ALIGN.CENTER)

    # Bottom text
    tb = add_textbox(slide, Inches(1.0), Inches(5.8), Inches(11), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, "Connected to the ", FONT_BODY, Pt(17), LIGHT)
    add_run(p, "Healthy Futures Initiative", FONT_BODY, Pt(17), GOLD, bold=True)
    add_run(p, ", building pathways for underserved communities into health professions", FONT_BODY, Pt(17), LIGHT)

    add_slide_number(slide, 13)


def slide_14_whats_next(prs):
    """Spring 2026 & Beyond"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BLUE)

    tb = add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5.5), Inches(0.4))
    set_text(tb.text_frame, "Looking Ahead", size=Pt(20), color=GOLD, bold=True)

    tb = add_textbox(slide, Inches(0.8), Inches(1.0), Inches(5.5), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_run(p, "Spring 2026\n& ", FONT_TITLE, Pt(36), WHITE, bold=True)
    add_run(p, "Beyond", FONT_TITLE, Pt(36), GOLD, bold=True)

    add_accent_line(slide, Inches(0.8), Inches(2.3))

    items = [
        ("Simulation Field Trip", "Grunin Center hands-on experience in Nursing, PA, OT, AT, PT"),
        ("Peer Mentorship Training", "Building public speaking & mentorship skills"),
        ("Future Leaders Symposium", "Full day mini-conference with certificates"),
        ("Community Outreach", "Connecting with local schools through Healthy Futures"),
    ]
    tb = add_textbox(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (bold_part, rest) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(10)
        add_run(p, "● ", FONT_BODY, Pt(17), GOLD)
        add_run(p, bold_part, FONT_BODY, Pt(17), WHITE, bold=True)
        add_run(p, f" — {rest}", FONT_BODY, Pt(17), LIGHT)

    # Right side: growth trajectory as text steps
    steps = [
        ("Program Design", RGBColor(0x1E, 0x3A, 0x5F)),
        ("Ice Breakers", RGBColor(0x1E, 0x6A, 0x5F)),
        ("Workshops", TEAL),
        ("Field Trips & Surveys", RGBColor(0x8A, 0xB3, 0x4F)),
        ("Symposium & Outreach", GOLD),
    ]
    for i, (step_text, color) in enumerate(steps):
        x = Inches(7.5)
        y = Inches(1.0 + i * 1.2)
        # Step circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.05), Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.color.rgb = WHITE
        circle.line.width = Pt(2)
        # Connector line (not for last)
        if i < len(steps) - 1:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.17), y + Inches(0.5), Pt(3), Inches(0.75))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(0x40, 0x50, 0x60)
            line.line.fill.background()
        # Label
        tb = add_textbox(slide, x + Inches(0.6), y, Inches(4), Inches(0.5))
        set_text(tb.text_frame, step_text, FONT_BODY, Pt(17), color, bold=True)

    tb = add_textbox(slide, Inches(7.5), Inches(6.8), Inches(4.5), Inches(0.3))
    set_text(tb.text_frame, "Program Growth Trajectory", size=Pt(13), color=DIM, italic=True, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 14)


def slide_15_lessons(prs):
    """Lessons Learned"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    tb = add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.4))
    set_text(tb.text_frame, "Lessons Learned", size=Pt(20), color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, "What Building This Program ", FONT_TITLE, Pt(34), WHITE, bold=True)
    add_run(p, "Taught Me", FONT_TITLE, Pt(34), GOLD, bold=True)

    lessons = [
        ("Start With People", "Great programs aren't built on paper — they're built on relationships. Listen to your students first, design second.", GOLD),
        ("Iterate, Don't Perfect", "The first workshop wasn't perfect. The second was better. Progress beats perfection every time.", TEAL),
        ("Leadership Is Learned By Doing", "I didn't just teach leadership — I had to practice it. Building this program was my own leadership lab.", CORAL),
    ]
    for i, (title, desc, color) in enumerate(lessons):
        x = Inches(0.6 + i * 4.1)
        card = add_card_bg(slide, x, Inches(1.8), Inches(3.8), Inches(4.8), color)
        # Title
        tb = add_textbox(slide, x + Inches(0.2), Inches(2.2), Inches(3.4), Inches(0.6))
        set_text(tb.text_frame, title, FONT_BODY, Pt(20), color, bold=True, alignment=PP_ALIGN.CENTER)
        # Description
        tb = add_textbox(slide, x + Inches(0.2), Inches(3.0), Inches(3.4), Inches(3.0))
        set_text(tb.text_frame, desc, size=Pt(16), color=LIGHT, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 15)


def slide_16_thankyou(prs):
    """Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    tb = add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.4))
    set_text(tb.text_frame, "Gratitude", size=Pt(20), color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, "Thank ", FONT_TITLE, Pt(38), WHITE, bold=True)
    add_run(p, "You", FONT_TITLE, Pt(38), GOLD, bold=True)

    thanks = [
        ("🎓", "Dr. Clifford", "Acting Dean & Mentor\nFor believing in the vision and guiding me every step", GOLD),
        ("🏫", "SNHS Faculty & Staff", "For the resources, support, and trust to build something new", TEAL),
        ("🤝", "20 Student Ambassadors", "For showing up, leaning in, and making this program come alive", CORAL),
        ("🏆", "Scholarship Week Committee", "For the opportunity to share this journey with you today", BLUE_ACC),
    ]
    for i, (icon, name, desc, color) in enumerate(thanks):
        x = Inches(0.5 + i * 3.15)
        card = add_card_bg(slide, x, Inches(1.8), Inches(2.9), Inches(4.2), color)
        # Icon
        tb = add_textbox(slide, x, Inches(2.0), Inches(2.9), Inches(0.6))
        set_text(tb.text_frame, icon, size=Pt(32), alignment=PP_ALIGN.CENTER)
        # Name
        tb = add_textbox(slide, x + Inches(0.1), Inches(2.7), Inches(2.7), Inches(0.5))
        set_text(tb.text_frame, name, FONT_BODY, Pt(17), color, bold=True, alignment=PP_ALIGN.CENTER)
        # Description
        tb = add_textbox(slide, x + Inches(0.1), Inches(3.3), Inches(2.7), Inches(2.2))
        set_text(tb.text_frame, desc, size=Pt(14), color=LIGHT, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(1.0), Inches(6.3), Inches(11), Inches(0.5))
    set_text(tb.text_frame, "None of this would have been possible without each of you", size=Pt(17), color=DIM, italic=True, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 16)


def slide_17_closing(prs):
    """Closing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, TEAL, GOLD)

    tb = add_textbox(slide, Inches(2.0), Inches(1.5), Inches(9.3), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, "The best way to learn\nleadership is to ", FONT_TITLE, Pt(36), WHITE, bold=True)
    add_run(p, "build something", FONT_TITLE, Pt(36), GOLD, bold=True)
    add_run(p, "\nthat matters.", FONT_TITLE, Pt(36), WHITE, bold=True)

    add_accent_line(slide, Inches(5.8), Inches(4.2), Inches(1.5))

    tb = add_textbox(slide, Inches(2.0), Inches(4.6), Inches(9.3), Inches(0.5))
    set_text(tb.text_frame, "Thank you for your time and support.", size=Pt(17), color=LIGHT, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(2.0), Inches(5.3), Inches(9.3), Inches(0.4))
    set_text(tb.text_frame, "Bingjun Li, M.S.Ed.", size=Pt(18), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(2.0), Inches(5.7), Inches(9.3), Inches(0.4))
    set_text(tb.text_frame, "Graduate Assistant · SNHS · Monmouth University", size=Pt(14), color=DIM, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(2.0), Inches(6.1), Inches(9.3), Inches(0.4))
    set_text(tb.text_frame, "Mentored by Dr. Clifford, Acting Dean", size=Pt(14), color=DIM, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 17)


# ─── Main ────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_hook(prs)
    slide_03_who_am_i(prs)
    slide_04_mission(prs)
    slide_05_what_i_built(prs)
    slide_06_timeline(prs)
    slide_07_workshop1(prs)
    slide_08_evolution(prs)
    slide_09_workshop2(prs)
    slide_10_feedback(prs)
    slide_11_challenges(prs)
    slide_12_breakthrough(prs)
    slide_13_impact(prs)
    slide_14_whats_next(prs)
    slide_15_lessons(prs)
    slide_16_thankyou(prs)
    slide_17_closing(prs)

    out_path = os.path.join(BASE, "presentation.pptx")
    prs.save(out_path)
    print(f"✅ Saved {out_path}")
    print(f"   {len(prs.slides)} slides generated")


if __name__ == "__main__":
    main()
